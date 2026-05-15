"""Generate 7-slide PPTX for RoadSoS hackathon submission."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG_DARK = RGBColor(10, 10, 10)
RED = RGBColor(220, 38, 38)
WHITE = RGBColor(248, 250, 252)
GRAY = RGBColor(148, 163, 184)
DARK_GRAY = RGBColor(71, 85, 105)
GREEN = RGBColor(22, 163, 74)
ORANGE = RGBColor(245, 158, 11)
BLUE = RGBColor(59, 130, 246)


def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return p


# ── SLIDE 1: Welcome ──
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide, RGBColor(15, 10, 10))

add_text(slide, 0, 1.5, 13.333, 1.2, "RoadSoS", 72, RED, True, PP_ALIGN.CENTER)
add_text(slide, 0, 2.8, 13.333, 0.8, "The AI Companion that turns every bystander into a first responder", 24, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, 0, 4.0, 13.333, 0.8, "Saving lives in the Golden Hour through AI-driven triage,\nlocation-aware emergency dispatch, and voice-first accessibility", 16, DARK_GRAY, False, PP_ALIGN.CENTER)
add_text(slide, 0, 5.5, 13.333, 0.5, "CoERS, IIT Madras  |  Road Safety Hackathon 2026  |  Team RBG Labs", 14, DARK_GRAY, False, PP_ALIGN.CENTER)
add_text(slide, 11.5, 7.0, 1.5, 0.4, "1 / 7", 11, DARK_GRAY, False, PP_ALIGN.RIGHT)

# ── SLIDE 2: The Problem ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0, 0.8, 13.333, 1.5, "1.19M", 96, RED, True, PP_ALIGN.CENTER)
add_text(slide, 0, 2.3, 13.333, 0.6, "lives lost on roads every year \u2014 half within 60 minutes", 20, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, 0, 3.2, 13.333, 0.8, '"Most aren\'t killed by the crash. They\'re killed by the wait."', 24, RGBColor(252, 165, 165), False, PP_ALIGN.CENTER)

tf = add_text(slide, 2, 4.3, 9.333, 2.5, "", 16, GRAY)
rows = [
    ("Metric", "Global", "BIMSTEC"),
    ("Ambulance response (urban)", "8\u201312 min", "20\u201335 min"),
    ("Ambulance response (rural)", "15\u201325 min", "45\u201390+ min"),
    ("Bystander first-aid rate", "15\u201325%", "< 5%"),
    ("Deaths within Golden Hour", "~50%", "60\u201370%"),
]
for i, (metric, glob, bimstec) in enumerate(rows):
    color = DARK_GRAY if i == 0 else GRAY
    bcolor = DARK_GRAY if i == 0 else RED
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"{metric:<35} {glob:<15} {bimstec}"
    p.font.size = Pt(14)
    p.font.color.rgb = color
    p.font.bold = (i == 0)

add_text(slide, 11.5, 7.0, 1.5, 0.4, "2 / 7", 11, DARK_GRAY, False, PP_ALIGN.RIGHT)

# ── SLIDE 3: Our Solution ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0, 0.5, 13.333, 0.8, "Our Solution", 48, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, 0, 1.4, 13.333, 0.7, "A multi-modal AI emergency companion that orchestrates the entire chain of survival", 20, GRAY, False, PP_ALIGN.CENTER)

features = [
    ("\u26a1 Speed", "One tap, zero friction.\nNo login walls."),
    ("\U0001f5e3 Vernacular", "8+ BIMSTEC languages.\nCode-mixed speech."),
    ("\U0001f4f6 Resilient", "Offline-first. SMS fallback.\nOne bar of signal is enough."),
]
for i, (title, desc) in enumerate(features):
    x = 1.5 + i * 3.8
    add_text(slide, x, 2.5, 3.2, 0.5, title, 22, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, x, 3.1, 3.2, 1.0, desc, 14, GRAY, False, PP_ALIGN.CENTER)

channels = "PWA  \u00b7  WhatsApp  \u00b7  IVR Voice  \u00b7  SMS  \u00b7  AR Overlay"
add_text(slide, 0, 4.8, 13.333, 0.5, channels, 16, DARK_GRAY, False, PP_ALIGN.CENTER)

add_text(slide, 11.5, 7.0, 1.5, 0.4, "3 / 7", 11, DARK_GRAY, False, PP_ALIGN.RIGHT)

# ── SLIDE 4: How It Works ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0, 0.3, 13.333, 0.8, "How It Works", 48, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, 0, 1.1, 13.333, 0.5, "11 PM, NH-37 Assam. A passing driver, Ravi, stops at a crash scene.", 18, GRAY, False, PP_ALIGN.CENTER)

timeline = [
    ("0:00", "\u2705", "Ravi taps SOS on PWA \u2014 or sends \"SOS\" on WhatsApp"),
    ("0:05", "\U0001f916", "GPS locked. AI asks triage questions in Assamese"),
    ("0:20", "\U0001f691", "HIGH severity \u2192 Ambulance, hospital, police dispatched"),
    ("0:35", "\U0001f4cb", "Voice + AR guides Ravi through bleeding control"),
    ("1:10", "\U0001f4f7", "Scene photos uploaded. ICE contacts notified"),
    ("11:20", "\U0001f3e5", "Ambulance arrives. ER pre-warned. Golden Hour intact."),
]
for i, (time, icon, text) in enumerate(timeline):
    y = 1.9 + i * 0.65
    add_text(slide, 2, y, 1.0, 0.5, icon, 20, WHITE, False, PP_ALIGN.CENTER)
    add_text(slide, 3, y, 1.2, 0.5, time, 14, ORANGE, True, PP_ALIGN.LEFT)
    add_text(slide, 4.2, y, 7, 0.5, text, 15, GRAY, False, PP_ALIGN.LEFT)

tech_pills = "LangGraph  \u00b7  Llama 3.1  \u00b7  FastAPI  \u00b7  Whisper ASR  \u00b7  PostgreSQL  \u00b7  WhatsApp API"
add_text(slide, 0, 6.2, 13.333, 0.5, tech_pills, 13, DARK_GRAY, False, PP_ALIGN.CENTER)
add_text(slide, 11.5, 7.0, 1.5, 0.4, "4 / 7", 11, DARK_GRAY, False, PP_ALIGN.RIGHT)

# ── SLIDE 5: Tech & AI Depth ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text(slide, 0, 0.3, 13.333, 0.8, "Tech & AI Depth", 48, WHITE, True, PP_ALIGN.CENTER)

layers = [
    ("PRESENTATION", "PWA \u00b7 WhatsApp \u00b7 IVR \u00b7 SMS \u00b7 AR Camera", RGBColor(220, 38, 38)),
    ("AI ORCHESTRATION", "LangGraph Agent \u00b7 RAG (WHO Protocols) \u00b7 ASR/TTS \u00b7 IndicTrans2", RGBColor(139, 92, 246)),
    ("SERVICES", "Hospitals \u00b7 Ambulances \u00b7 Police \u00b7 Towing \u00b7 Puncture Shops", BLUE),
    ("DATA", "PostgreSQL \u00b7 pgvector \u00b7 Redis \u00b7 Celery", GREEN),
]
for i, (label, desc, color) in enumerate(layers):
    y = 1.5 + i * 1.1
    add_text(slide, 1, y, 6.5, 0.4, label, 16, color, True, PP_ALIGN.CENTER)
    add_text(slide, 1, y + 0.4, 6.5, 0.4, desc, 13, GRAY, False, PP_ALIGN.CENTER)

# Right side callout
tf = add_text(slide, 8, 1.8, 4.5, 2.0, "Why RAG, not free generation?", 18, RGBColor(252, 165, 165), True)
add_para(tf, "", 6, GRAY)
add_para(tf, "Medical hallucination is unacceptable. Every instruction grounded in vetted WHO + ATLS corpus.", 14, GRAY)
add_para(tf, "", 6, GRAY)
add_para(tf, "Confidence < 0.78 \u2192 human dispatcher", 14, ORANGE, True)

pills2 = "Open-weights LLM  \u00b7  On-prem deploy  \u00b7  Offline-first  \u00b7  8 languages"
add_text(slide, 8, 4.5, 4.5, 0.5, pills2, 12, DARK_GRAY, False, PP_ALIGN.LEFT)
add_text(slide, 11.5, 7.0, 1.5, 0.4, "5 / 7", 11, DARK_GRAY, False, PP_ALIGN.RIGHT)

# ── SLIDE 6: Impact & Roadmap ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, RGBColor(8, 12, 8))

add_text(slide, 0, 0.3, 13.333, 0.8, "Impact & Roadmap", 48, WHITE, True, PP_ALIGN.CENTER)

# Left: impact stat
add_text(slide, 1, 1.8, 5, 1.2, "22\u201334%", 72, GREEN, True, PP_ALIGN.CENTER)
add_text(slide, 1, 3.2, 5, 0.5, "projected mortality reduction", 18, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, 1, 3.8, 5, 1.0, "Modelled against WHO trauma curves.\nIf response time cut by 15 minutes\non BIMSTEC highway corridors.", 13, DARK_GRAY, False, PP_ALIGN.CENTER)

# Right: roadmap
roadmap = [
    ("Stage 1", "Jun 2026", "PWA MVP, WhatsApp bot, TN seed DB, demo"),
    ("Pilot A", "Q3 2026", "2 highway corridors (NH-44, NH-66) with state EMS"),
    ("Pilot B", "Q4 2026", "6 Indian states + Sri Lanka + Bhutan via BIMSTEC"),
    ("Scale", "2027", "Cross-BIMSTEC rollout, national 112 integration"),
]
for i, (phase, date, desc) in enumerate(roadmap):
    y = 1.8 + i * 0.9
    add_text(slide, 7, y, 1.5, 0.5, phase, 15, ORANGE, True)
    add_text(slide, 8.5, y, 1.5, 0.5, date, 13, DARK_GRAY, False)
    add_text(slide, 10, y, 3, 0.5, desc, 13, GRAY, False)

add_text(slide, 11.5, 7.0, 1.5, 0.4, "6 / 7", 11, DARK_GRAY, False, PP_ALIGN.RIGHT)

# ── SLIDE 7: Thank You ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, RGBColor(15, 10, 10))

add_text(slide, 0, 1.0, 13.333, 1.0, '"Every minute matters.\nLet\'s give every road user a fighting chance."', 28, RGBColor(251, 191, 36), False, PP_ALIGN.CENTER)

# QR codes
qr_dir = os.path.dirname(os.path.abspath(__file__))
qr_pwa = os.path.join(qr_dir, "qr-pwa.png")
qr_api = os.path.join(qr_dir, "qr-api.png")

if os.path.exists(qr_pwa):
    slide.shapes.add_picture(qr_pwa, Inches(4.5), Inches(3.0), Inches(1.5), Inches(1.5))
    add_text(slide, 4.5, 4.6, 1.5, 0.4, "Try the PWA", 11, DARK_GRAY, False, PP_ALIGN.CENTER)

if os.path.exists(qr_api):
    slide.shapes.add_picture(qr_api, Inches(7.3), Inches(3.0), Inches(1.5), Inches(1.5))
    add_text(slide, 7.3, 4.6, 1.5, 0.4, "API Documentation", 11, DARK_GRAY, False, PP_ALIGN.CENTER)

add_text(slide, 0, 5.5, 13.333, 0.5, "Team RoadSoS  \u00b7  CoERS, IIT Madras  \u00b7  RBG Labs", 16, GRAY, False, PP_ALIGN.CENTER)
add_text(slide, 0, 6.1, 13.333, 0.5, "Road Safety Hackathon 2026  \u00b7  BIMSTEC Track", 13, DARK_GRAY, False, PP_ALIGN.CENTER)

# Live URLs
add_text(slide, 0, 6.6, 13.333, 0.4, "PWA: frontend-ten-gamma-78.vercel.app  |  API: rare-adventure-production-984a.up.railway.app/docs  |  Code: github.com/leonkaushikdeka/roadsos", 10, DARK_GRAY, False, PP_ALIGN.CENTER)
add_text(slide, 11.5, 7.0, 1.5, 0.4, "7 / 7", 11, DARK_GRAY, False, PP_ALIGN.RIGHT)

# SAVE
out = os.path.join(qr_dir, "RoadSoS_Presentation.pptx")
prs.save(out)
print(f"Saved: {out}")
